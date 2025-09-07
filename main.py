
from fastapi import FastAPI, UploadFile
import google.generativeai as genai
import PyPDF2
import docx  
from pptx import Presentation  
import os
from dotenv import load_dotenv
from fastapi.middleware.cors import CORSMiddleware
import json

load_dotenv()

app = FastAPI()

website_url = os.getenv("https://revision-tool-mu.vercel.app/", "http://localhost:3000")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[website_url],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)




api_key = os.getenv("GEMINI_API_KEY")
if not api_key:
    raise ValueError("⚠️ La clé API Gemini n'est pas définie dans le fichier .env")

genai.configure(api_key=api_key)
model = genai.GenerativeModel("gemini-1.5-flash")


def extract_text_from_pdf(file):
    reader = PyPDF2.PdfReader(file)
    pages_content = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if text:
            pages_content.append({"page": i, "text": text.strip()})
    return pages_content


def extract_text_from_docx(file):
    doc = docx.Document(file)
    full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    return [{"page": None, "text": full_text}]  


def extract_text_from_pptx(file):
    prs = Presentation(file)
    slides_content = []
    for i, slide in enumerate(prs.slides, start=1):
        text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()])
        if text:
            slides_content.append({"page": i, "text": text.strip()})
    return slides_content


@app.post("/analyser/")
async def analyser_doc(file: UploadFile):
    """Endpoint qui supporte PDF, DOCX, DOC, PPTX"""
    filename = file.filename.lower()

    if filename.endswith(".pdf"):
        pages_content = extract_text_from_pdf(file.file)
    elif filename.endswith(".docx") or filename.endswith(".doc"):
        pages_content = extract_text_from_docx(file.file)
    elif filename.endswith(".pptx") or filename.endswith(".ppt"):
        pages_content = extract_text_from_pptx(file.file)
    else:
        return {"error": "Format de fichier non supporté"}

    # Construire le document avec pages
    document_str = "\n\n".join(
        [f"[PAGE {p['page']}]\n{p['text']}" if p['page'] else p['text'] for p in pages_content]
    )

    prompt = f"""
Tu es un assistant pédagogique expert. 
Je vais te fournir un document complet, page par page.
Ta tâche est de générer trois sections distinctes à partir du document :

1. "resume" : Résumé clair et concis du document (<p>)
2. "questions" : Questions sous deux catégories, numérotées, avec HTML
    - "generales" : questions de type général (1., 2., 3., ...)
    - "detaillees" : questions plus spécifiques (1., 2., 3., ...)
3. "fiche" : Tableau JSON de parties de révision :
    - Chaque objet du tableau doit avoir :
        - "titre": titre de la partie
        - "resume": résumé de la partie (<p>)
        - "points": liste des points clés, numérotés de la même façon que les questions (1., 2., 3., ...)
        - "numero_page": numéro de la page correspondant

⚠️ Réponds uniquement en JSON valide avec ces clés :
{{
  "resume": "<p>Résumé ici</p>",
  "questions": {{
       "generales": "1. Question 1 ...\\n2. Question 2 ...",
       "detaillees": "1. Question détaillée 1 ...\\n2. Question détaillée 2 ..."
  }},
  "fiche": [
      {{
        "titre": "Titre de la partie",
        "resume": "<p>Résumé de la partie</p>",
        "points": "1. Point 1 ...\\n2. Point 2 ...\\n3. Point 3 ...",
        "numero_page": 5
      }}
  ]
}}

Document (page par page) :
{document_str}
"""

    response = model.generate_content(prompt)

    try:
        start = response.text.find("{")
        end = response.text.rfind("}") + 1
        result = json.loads(response.text[start:end])
    except json.JSONDecodeError:
        result = {
            "resume": "",
            "questions": {"generales": "", "detaillees": ""},
            "fiche": []
        }

    return result


@app.post("/generatequiz/")
async def generatequiz(file: UploadFile):
    """Génère un quiz à partir du document uploadé avec 3 options par question"""
    filename = file.filename.lower()

    if filename.endswith(".pdf"):
        pages_content = extract_text_from_pdf(file.file)
    elif filename.endswith(".docx") or filename.endswith(".doc"):
        pages_content = extract_text_from_docx(file.file)
    elif filename.endswith(".pptx") or filename.endswith(".ppt"):
        pages_content = extract_text_from_pptx(file.file)
    else:
        return {"error": "Format de fichier non supporté"}

    document_str = "\n\n".join(
        [f"[PAGE {p['page']}]\n{p['text']}" if p['page'] else p['text'] for p in pages_content]
    )

    prompt = f"""
Tu es un assistant pédagogique expert.
Génère un quiz à partir du document fourni (page par page). 
Pour chaque question, donne exactement 3 réponses : 1 correcte et 2 incorrectes.
Pour chaque réponse, écris une explication claire et précise :
- si la réponse est correcte, explique exactement pourquoi elle l'est en te basant sur le document,
- si la réponse est incorrecte, explique exactement pourquoi elle est fausse, et ce que l'on doit retenir.
Ne te contente pas de dire "voir page X", sois spécifique.
Répond uniquement en JSON comme ceci :
{{
  "questions": [
    {{
      "question": "Texte de la question",
      "reponses": [
        {{"texte": "Réponse 1", "correct": true/false, "explication": "Explication claire"}}
      ]
    }}
  ]
}}
{document_str}
"""


    response = model.generate_content(prompt)

    try:
        start = response.text.find("{")
        end = response.text.rfind("}") + 1
        quiz = json.loads(response.text[start:end])
    except json.JSONDecodeError:
        quiz = {"questions": []}

    return quiz
